from pptx import Presentation
from pptx.util import Inches
import os
import re



class Pptx:
	def __init__(self, model):
		self.prs = Presentation(model)


	def PopulatePageByName(self, name):
		templateStyleNum = len(self.prs.slide_layouts)
		oneSlide = self.prs.slides.add_slide(self.prs.slide_layouts[0])
		body_shapes = oneSlide.shapes.placeholders

		for body_shape in body_shapes:
		    body_shape.text = name
		    break

	def Save(self, output_name):
		self.prs.save(output_name)

def UnifyName():
		basepath = os.getcwd() + '/slides/'

		for filename in os.listdir("./slides"):
			digits = re.sub('\D', '', filename)
			digits = int(digits)
			os.rename(basepath + filename, basepath + str(digits) + ".png")

# def main():
	# UnifyName()
	# pptx_writer = Pptx("model.pptx")
# 	pptx_writer.PopulatePageByName("what we like")
# 	pptx_writer.PopulatePageByName("what we love")
# 	pptx_writer.Save("one.pptx")

if __name__ == "__main__":
    main()